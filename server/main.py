"""
VDT Slack Bot Server
/vdt CVS-XXXXX 슬래시 커맨드 → 가상 개발팀 기획서 품질 분석 → #qa-ai-report 전송
"""
import json, os, re, threading, tempfile, subprocess, datetime
from pathlib import Path
from concurrent.futures import ThreadPoolExecutor

from slack_bolt import App
from slack_bolt.adapter.flask import SlackRequestHandler
from flask import Flask, request
import anthropic

# ── 크레덴셜 로드 ──
def _env_or_file(env_key, file_path, json_key):
    val = os.environ.get(env_key)
    if val:
        return val
    p = Path.home() / ".bagelcode" / file_path
    if p.exists():
        return json.load(open(p))[json_key]
    raise RuntimeError(f"크레덴셜 없음: {env_key} 또는 ~/.bagelcode/{file_path}")

SLACK_TOKEN   = _env_or_file("SLACK_BOT_TOKEN",     "slack_bot.json", "bot_token")
SLACK_SIGN    = _env_or_file("SLACK_SIGNING_SECRET", "slack_bot.json", "signing_secret")
JIRA_DOMAIN_V = _env_or_file("JIRA_DOMAIN",          "jira.json",      "domain")
JIRA_EMAIL_V  = _env_or_file("JIRA_EMAIL",            "jira.json",      "email")
JIRA_TOKEN_V  = _env_or_file("JIRA_TOKEN",            "jira.json",      "token")
AIPROXY_TOKEN = _env_or_file("ANTHROPIC_API_KEY",     "anthropic.json", "api_key")
GITHUB_TOKEN_V = os.environ.get("GITHUB_TOKEN", "")
VIEWER_PATH   = os.environ.get("VIEWER_PATH", "/app/viewer.html")
CHANNEL       = os.environ.get("REPORT_CHANNEL", "C0AQTSRRFHC")
PAGES_BASE    = os.environ.get("PAGES_BASE_URL", "https://bumsinkim-cyber.github.io/vdt-reports")
PAGES_REPO    = os.environ.get("PAGES_REPO", "bumsinkim-cyber/vdt-reports")
REPO_DIR      = "/tmp/vdt-reports-git"

JIRA_BASE = f"https://{JIRA_DOMAIN_V}"
JIRA_AUTH = (JIRA_EMAIL_V, JIRA_TOKEN_V)

app       = App(token=SLACK_TOKEN, signing_secret=SLACK_SIGN)
flask_app = Flask(__name__)
handler   = SlackRequestHandler(app)

claude = anthropic.Anthropic(
    api_key=AIPROXY_TOKEN,
    base_url=os.environ.get("ANTHROPIC_BASE_URL", "https://api.anthropic.com"),
    timeout=300.0,
)

# ── 슬래시 커맨드 ──
@app.command("/vdt")
def handle_vdt(ack, body, say):
    ack()
    ticket_key = (body.get("text") or "").strip().upper()
    if not re.match(r"^[A-Z]+-\d+$", ticket_key):
        say("❌ 올바른 티켓 키를 입력해 주세요. 예: `/vdt CVS-10785`")
        return
    app.client.chat_postMessage(
        channel=CHANNEL,
        text=f"🔍 *{ticket_key}* 분석을 시작합니다. 잠시만 기다려 주세요..."
    )
    threading.Thread(target=run_vdt, args=(ticket_key,), daemon=True).start()


def run_vdt(ticket_key: str):
    try:
        tmpdir = Path(tempfile.mkdtemp(prefix=f"vdt-{ticket_key}-"))
        _run_pipeline(ticket_key, tmpdir)
    except Exception as e:
        import traceback
        app.client.chat_postMessage(
            channel=CHANNEL,
            text=f"❌ *{ticket_key}* 분석 중 오류가 발생했습니다.\n```{traceback.format_exc()[-800:]}```"
        )


def _notify(ticket_key: str, msg: str):
    app.client.chat_postMessage(channel=CHANNEL, text=f"[{ticket_key}] {msg}")


# ── 파이프라인 ──
def _run_pipeline(ticket_key: str, tmpdir: Path):
    import requests as req

    # STEP 1: Jira
    _notify(ticket_key, "📋 Jira 티켓 수집 중...")
    ticket_res = req.get(
        f"{JIRA_BASE}/rest/api/2/issue/{ticket_key}",
        params={"fields": "summary,description,subtasks,assignee,status"},
        auth=JIRA_AUTH,
    )
    ticket = ticket_res.json()
    fields = ticket.get("fields", {})
    title  = fields.get("summary", ticket_key)
    subtasks = [
        {"key": s["key"], "summary": s["fields"].get("summary", ""),
         "status": s["fields"].get("status", {}).get("name", "")}
        for s in fields.get("subtasks", [])
    ]
    rl_res = req.get(f"{JIRA_BASE}/rest/api/2/issue/{ticket_key}/remotelink", auth=JIRA_AUTH)
    remotelinks = rl_res.json()

    # STEP 2: 기획서 수집
    _notify(ticket_key, "📄 기획서 수집 중...")
    spec_text = _collect_specs(tmpdir, remotelinks)
    _notify(ticket_key, f"✅ 기획서 수집 완료 ({len(spec_text)}자). 에이전트 분석 시작...")

    ticket_info = json.dumps({"title": title, "assignee": ""}, ensure_ascii=False)
    subtasks_json = json.dumps(subtasks, ensure_ascii=False)

    # STEP 3: Planner
    _notify(ticket_key, "🗂 Planner 분석 중...")
    planner_json = _call_agent(_planner_prompt(spec_text, ticket_info, ticket_key, subtasks_json))

    # STEP 4: Developer + Artist 병렬
    _notify(ticket_key, "💻🎨 Developer + Artist 동시 분석 중...")
    with ThreadPoolExecutor(max_workers=2) as ex:
        dev_fut = ex.submit(_call_agent, _developer_prompt(planner_json))
        art_fut = ex.submit(_call_agent, _artist_prompt(planner_json))
        developer_json = dev_fut.result()
        artist_json    = art_fut.result()

    # STEP 4.5: TA
    _notify(ticket_key, "⚙️ TA 분석 중...")
    ta_json = _call_agent(_ta_prompt(planner_json, developer_json, artist_json))

    # STEP 5: QA Pre-analyst
    _notify(ticket_key, "🔍 QA 기획서 품질 분석 중...")
    qa_json = _call_agent(_qa_prompt(planner_json, developer_json, artist_json, ta_json), max_tokens=8192)

    # STEP 5.5: 팀 회의
    _notify(ticket_key, "💬 팀 회의 시뮬레이션 중...")
    meeting_json = _call_agent(_meeting_prompt(planner_json, developer_json, artist_json, ta_json, qa_json))

    for name, data in [("planner", planner_json), ("developer", developer_json),
                        ("artist", artist_json), ("ta", ta_json),
                        ("qa", qa_json), ("meeting", meeting_json)]:
        (tmpdir / f"{name}.json").write_text(data, encoding="utf-8")

    # STEP 6: HTML
    report_path = tmpdir / "report.html"
    _generate_html(ticket_key, tmpdir, report_path)

    # STEP 6.5: GitHub Pages 배포
    report_url = _deploy_to_pages(ticket_key, report_path, json.loads(planner_json))

    # STEP 7: Slack 전송
    _send_slack(ticket_key, title, planner_json, developer_json, qa_json, report_url)


# ── Google 인증 ──
def _get_google_creds():
    import base64
    from google.oauth2.credentials import Credentials
    from google.auth.transport.requests import Request

    SCOPES = [
        "https://www.googleapis.com/auth/drive.readonly",
        "https://www.googleapis.com/auth/presentations.readonly",
    ]
    token_b64 = os.environ.get("GOOGLE_TOKEN_JSON")
    if token_b64:
        info = json.loads(base64.b64decode(token_b64).decode())
    else:
        p = Path.home() / ".bagelcode" / "google_token.json"
        if not p.exists():
            return None
        info = json.loads(p.read_text())

    creds = Credentials.from_authorized_user_info(info, SCOPES)
    if creds.expired and creds.refresh_token:
        creds.refresh(Request())
    return creds


# ── 기획서 수집 ──
def _get_slides_text(file_id: str, creds) -> str:
    import io
    from googleapiclient.discovery import build
    from googleapiclient.http import MediaIoBaseDownload
    from pptx import Presentation

    drive = build("drive", "v3", credentials=creds)
    meta  = drive.files().get(fileId=file_id, fields="mimeType").execute()
    mime  = meta.get("mimeType", "")

    if mime == "application/vnd.google-apps.presentation":
        content = drive.files().export(fileId=file_id, mimeType="text/plain").execute()
        return (content.decode("utf-8", errors="ignore") if isinstance(content, bytes) else str(content)).strip()

    buf = io.BytesIO()
    req = drive.files().get_media(fileId=file_id)
    dl  = MediaIoBaseDownload(buf, req)
    done = False
    while not done:
        _, done = dl.next_chunk()

    prs = Presentation(buf)
    chunks = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(r.text for r in para.runs).strip()
                    if t:
                        texts.append(t)
        if texts:
            chunks.append(f"=== 슬라이드 {i} ===\n" + "\n".join(texts))
        else:
            chunks.append(f"=== 슬라이드 {i} === [이미지 전용 슬라이드 — 텍스트 추출 불가]")
    return "\n\n".join(chunks)


def _collect_specs(tmpdir: Path, remotelinks: list) -> str:
    import requests as req
    creds  = _get_google_creds()
    chunks = []

    for lk in remotelinks:
        url   = lk.get("object", {}).get("url", "")
        title = lk.get("object", {}).get("title", "")

        if "docs.google.com/presentation" in url or "drive.google.com" in url:
            m = re.search(r"/(?:presentation|file)/d/([a-zA-Z0-9_-]+)", url)
            if not m:
                m = re.search(r"[?&]id=([a-zA-Z0-9_-]+)", url)
            if m and creds:
                text = _get_slides_text(m.group(1), creds)
                if text:
                    chunks.append(f"=== Google Slides: {title} ===\n{text}")
            continue

        if "atlassian.net/wiki" in url:
            m = re.search(r"pageId=(\d+)", url)
            if not m:
                continue
            page_id = m.group(1)
            res = req.get(
                f"{JIRA_BASE}/wiki/rest/api/content/{page_id}",
                params={"expand": "body.storage"},
                auth=JIRA_AUTH,
            )
            if not res.ok:
                continue
            import html as html_mod
            body = res.json().get("body", {}).get("storage", {}).get("value", "")
            text = re.sub(r"<[^>]+>", " ", body)
            text = html_mod.unescape(re.sub(r"\s+", " ", text)).strip()

            # Confluence 본문이 Google Drive 파일 embed인 경우
            slides_m = re.search(r"https://docs\.google\.com/presentation/d/([a-zA-Z0-9_-]+)", text)
            if slides_m and creds:
                slides_text = _get_slides_text(slides_m.group(1), creds)
                if slides_text:
                    chunks.append(f"=== Google Slides (via Confluence: {title}) ===\n{slides_text}")
                    continue

            if text:
                chunks.append(f"=== Confluence: {title} ===\n{text}")

    return "\n\n---\n\n".join(chunks) if chunks else "[기획서 없음]"


# ── Claude 호출 ──
def _extract_json(text: str) -> str:
    start = text.find("{")
    if start == -1:
        return text
    depth = 0
    for i, c in enumerate(text[start:], start):
        if c == "{":
            depth += 1
        elif c == "}":
            depth -= 1
            if depth == 0:
                return text[start:i + 1]
    return text


def _call_agent(prompt: str, max_tokens: int = 4096) -> str:
    msg = claude.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=max_tokens,
        messages=[{"role": "user", "content": prompt}],
    )
    return _extract_json(msg.content[0].text)


# ── 에이전트 프롬프트 ──
def _planner_prompt(spec_text, ticket_info, ticket_key, subtasks_json):
    return f"""[Planner 역할 지시]
당신은 시니어 기획자입니다. 아래 기획서 원문을 읽고 JSON을 출력하라.

원문 충실도 규칙: 의미 단위 재해석 금지 / 누락 요구사항 보완 금지 / 원문에 없는 섹션 생성 금지
[이미지 전용 슬라이드] 항목은 ac_items 제외, image_only_slides 배열에 슬라이드 번호 추가

AC 카테고리: M=메인/매칭/메타 U=UI/UX S=시스템/서버 A=애니메이션/아트 L=로직 E=예외/엣지케이스

티켓 정보: {ticket_info}
서브태스크: {subtasks_json}

기획서 원문:
---
{spec_text[:8000]}
---

JSON 코드 블록으로만 출력:
{{"title":"","ticket_key":"{ticket_key}","purpose":"","key_features":[{{"component":"","content":""}}],"subtasks":[],"incomplete_areas":[{{"comment":""}}],"ui_ux_changes":[],"edge_cases":[],"image_only_slides":[],"ac_items":[{{"id":"M-01","category":"M","category_name":"메인/매칭","content":"","source":""}}]}}"""


def _developer_prompt(planner_json):
    return f"""[Developer 역할 지시]
당신은 시니어 개발자(Tech Lead)입니다. Planner 분석 결과를 바탕으로 개발 리스크와 구현 계획을 JSON으로 출력하라.

Planner 결과: {planner_json}

JSON 코드 블록으로만 출력:
{{"impact_scope":[{{"module":"","change_type":"신규|수정|공용 수정|삭제|완료","is_risky":true,"detail":""}}],"impl_risks":[{{"feature":"","complexity":"High|Med|Low","risk":""}}],"impl_order":[""],"tech_notes":[""]}}"""


def _artist_prompt(planner_json):
    return f"""[Artist 역할 지시]
당신은 시니어 게임 아티스트/UI 디자이너입니다. Planner 분석 결과를 바탕으로 에셋 목록과 UI/UX 플로우를 JSON으로 출력하라.

Planner 결과: {planner_json}

JSON 코드 블록으로만 출력:
{{"ui_ux_flow":["화면A → (액션) → 화면B"],"assets":{{"images":[{{"name":"","usage":"","size":"","is_new":true}}],"animations":[{{"name":"","type":"","trigger":""}}],"effects":[{{"name":"","trigger":"","note":""}}],"fonts":[]}},"design_guidelines":[],"total_asset_days":0,"asset_summary":""}}"""


def _ta_prompt(planner_json, developer_json, artist_json):
    return f"""[TA 역할 지시]
당신은 테크니컬 아티스트입니다. 아트 요구사항과 개발 계획을 바탕으로 기술 스펙과 리스크를 JSON으로 출력하라.

Planner: {planner_json[:2000]}
Developer: {developer_json[:1000]}
Artist: {artist_json[:2000]}

JSON 코드 블록으로만 출력:
{{"asset_specs":[{{"name":"","format":"PNG","resolution":"","atlas":"Y/N","memory_est":""}}],"effect_specs":[{{"name":"","impl_method":"","perf_note":""}}],"pipeline_impact":{{"build_changed":false,"import_changes":[],"bundle_note":""}},"performance_notes":[{{"item":"","impact":"","mitigation":""}}],"ta_risks":[{{"item":"","risk":"","recommendation":""}}]}}"""


def _qa_prompt(planner_json, developer_json, artist_json, ta_json):
    planner  = json.loads(planner_json)
    ac_items = json.dumps(planner.get("ac_items", []), ensure_ascii=False)
    return f"""[QA Pre-analyst 역할 지시]
당신은 기획서 품질 분석가입니다. "개발자가 이 기획서만 보고 착수할 수 있는가"의 관점에서 각 AC 항목을 판정하고, ⚠️/🔴 항목에 대해 QA 테스트 예측을 분석하여 JSON으로 출력하라.

판정 기준:
✅ = 착수가능 (기획서만으로 구현 방향 결정 가능)
⚠️ = 확인필요 (세부 조건 불명확)
🔴 = 착수불가 (핵심 정보 누락 또는 상충)

각 항목 분석:
- ambiguity: 모호한 표현 (원문 인용, 없으면 빈 문자열)
- missing_info: 구현을 위해 없는 정보 (없으면 빈 문자열)
- question: 기획자에게 확인해야 할 질문 (없으면 빈 문자열)
- qa_prediction: ⚠️/🔴 판정 시, 이 모호함이 해소되지 않고 개발된 경우 QA에서 발생할 것으로 예측되는 구체적 문제. "어떤 케이스에서 → 어떤 증상" 구조로 기술. ✅이면 빈 문자열.

image_only_slides에 포함된 슬라이드 기반 항목은 분석 제외.
TA 리스크 항목이 있으면 해당 AC의 verdict를 ⚠️ 이상으로 올려라.

AC 항목: {ac_items}
개발 리스크: {developer_json[:1500]}
Artist: {artist_json[:1000]}
TA 리스크: {ta_json[:1000]}

JSON 코드 블록으로만 출력:
{{"spec_quality":[{{"id":"M-01","verdict":"✅|⚠️|🔴","verdict_label":"착수가능|확인필요|착수불가","ambiguity":"","missing_info":"","question":"","qa_prediction":""}}],"critical_questions":[{{"priority":"P0|P1","related_ids":["M-01"],"question":"","impact":""}}]}}"""


def _meeting_prompt(planner_json, developer_json, artist_json, ta_json, qa_json):
    planner   = json.loads(planner_json)
    developer = json.loads(developer_json)
    qa        = json.loads(qa_json)
    title     = planner.get("title", "")
    incomplete = " / ".join(i.get("comment", "") for i in planner.get("incomplete_areas", []))
    dev_risks  = " / ".join(f"{r['feature']}({r['complexity']})" for r in developer.get("impl_risks", [])[:5])
    qa_risks   = " / ".join(
        f"{a['id']}: {a.get('missing_info', '')}"
        for a in qa.get("spec_quality", []) if a.get("verdict") == "🔴"
    )
    artist_summary = json.loads(artist_json).get("asset_summary", "")
    ta_risks = " / ".join(f"{r['item']}: {r['risk']}" for r in json.loads(ta_json).get("ta_risks", [])[:3])
    cqs = " / ".join(
        f"[P0] {q['question'][:60]}"
        for q in qa.get("critical_questions", []) if q.get("priority") == "P0"
    )

    return f"""[팀 회의 퍼실리테이터]
기획→개발→아트→TA→QA 순서로 2라운드 팀 회의를 시뮬레이션하고 합의 액션을 도출하라.
강조할 내용은 **bold** 마크다운 사용. 합의 결과는 실행 가능한 액션 아이템으로 정리(P0=즉시, P1=스프린트 내, P2=권고).

기획서: {title}
미완성: {incomplete or '없음'}
🔴 착수불가 항목: {qa_risks or '없음'}
개발 리스크: {dev_risks or '없음'}
아트 요약: {artist_summary or '없음'}
TA 리스크: {ta_risks or '없음'}
Critical Q (P0): {cqs or '없음'}

JSON 코드 블록으로만 출력:
{{"date":"{datetime.date.today()}","participants":["기획 에이전트","개발 에이전트","아트 에이전트","TA 에이전트","QA 에이전트"],"rounds":[{{"round":1,"agent":"기획 에이전트","agent_type":"planner","content":""}}],"consensus":[{{"priority":"P0|P1|P2","action":""}}]}}"""


# ── HTML 보고서 생성 ──
def _generate_html(ticket_key: str, tmpdir: Path, output: Path):
    viewer   = Path(VIEWER_PATH).read_text(encoding="utf-8")
    planner  = json.loads((tmpdir / "planner.json").read_text())
    developer= json.loads((tmpdir / "developer.json").read_text())
    qa       = json.loads((tmpdir / "qa.json").read_text())
    artist   = json.loads((tmpdir / "artist.json").read_text()) if (tmpdir / "artist.json").exists() else {}
    ta       = json.loads((tmpdir / "ta.json").read_text())     if (tmpdir / "ta.json").exists()     else {}
    meeting  = json.loads((tmpdir / "meeting.json").read_text()) if (tmpdir / "meeting.json").exists() else {}

    ac_map = {ac["id"]: dict(ac) for ac in planner.get("ac_items", [])}
    for a in qa.get("spec_quality", []):
        if a["id"] in ac_map:
            ac_map[a["id"]].update(a)

    label_map = {"✅": "착수가능", "⚠️": "확인필요", "🔴": "착수불가"}
    stats = {v: 0 for v in label_map.values()}
    for ac in ac_map.values():
        v = ac.get("verdict")
        if v:
            stats[label_map.get(v, "확인필요")] += 1

    data = {
        "ticket_key":    ticket_key,
        "title":         planner.get("title", ticket_key),
        "analysis_date": datetime.date.today().strftime("%Y-%m-%d"),
        "status":        "분석 완료",
        "stats":         stats,
        "spec": {
            "purpose":          planner.get("purpose", ""),
            "key_features":     planner.get("key_features", []),
            "subtasks":         planner.get("subtasks", []),
            "incomplete_areas": planner.get("incomplete_areas", []),
            "ui_ux_changes":    planner.get("ui_ux_changes", []),
            "edge_cases":       planner.get("edge_cases", []),
        },
        "ac_items": list(ac_map.values()),
        "developer": {
            "impact_scope": developer.get("impact_scope", []),
            "impl_risks":   developer.get("impl_risks", []),
            "impl_order":   developer.get("impl_order", []),
            "tech_notes":   developer.get("tech_notes", []),
        },
        "artist": {
            "ui_ux_flow":        artist.get("ui_ux_flow", []),
            "assets":            artist.get("assets", {}),
            "design_guidelines": artist.get("design_guidelines", []),
            "total_asset_days":  artist.get("total_asset_days", 0),
            "asset_summary":     artist.get("asset_summary", ""),
        },
        "ta": {
            "asset_specs":       ta.get("asset_specs", []),
            "effect_specs":      ta.get("effect_specs", []),
            "pipeline_impact":   ta.get("pipeline_impact", {}),
            "performance_notes": ta.get("performance_notes", []),
            "ta_risks":          ta.get("ta_risks", []),
        },
        "critical_questions": qa.get("critical_questions", []),
        "meeting": meeting or None,
    }
    html = viewer.replace("__DATA_PLACEHOLDER__", json.dumps(data, ensure_ascii=False))
    output.write_text(html, encoding="utf-8")


# ── GitHub Pages 배포 ──
def _deploy_to_pages(ticket_key: str, report_path: Path, planner: dict) -> str:
    import shutil
    repo_dir = Path(REPO_DIR)
    gh_token = GITHUB_TOKEN_V
    repo_url = (
        f"https://{gh_token}@github.com/{PAGES_REPO}.git"
        if gh_token else f"https://github.com/{PAGES_REPO}.git"
    )

    if not (repo_dir / ".git").exists():
        subprocess.run(["git", "clone", repo_url, str(repo_dir)], check=True)
    else:
        if gh_token:
            subprocess.run(["git", "-C", str(repo_dir), "remote", "set-url", "origin", repo_url], check=True)
        subprocess.run(["git", "-C", str(repo_dir), "pull", "origin", "main"], check=True)

    subprocess.run(["git", "-C", str(repo_dir), "config", "user.email", "vdt-bot@bagelcode.com"], check=True)
    subprocess.run(["git", "-C", str(repo_dir), "config", "user.name", "VDT Bot"], check=True)

    filename = f"{ticket_key}.html"
    shutil.copy(report_path, repo_dir / filename)

    idx_path = repo_dir / "index.html"
    idx  = idx_path.read_text(encoding="utf-8")
    date = datetime.date.today().strftime("%Y-%m-%d")
    title = planner.get("title", ticket_key)
    new_item = f'  <li><a href="{filename}">{ticket_key} — {title}</a><div class="meta">{date}</div></li>'
    if filename not in idx:
        idx = idx.replace('<ul class="report-list" id="list">', f'<ul class="report-list" id="list">\n{new_item}')
        idx_path.write_text(idx, encoding="utf-8")

    subprocess.run(["git", "-C", str(repo_dir), "add", filename, "index.html"], check=True)
    subprocess.run(["git", "-C", str(repo_dir), "commit", "-m", f"report: {ticket_key} 분석 보고서 추가"], check=True)
    subprocess.run(["git", "-C", str(repo_dir), "push", "origin", "main"], check=True)

    return f"{PAGES_BASE}/{filename}"


# ── Slack 결과 전송 ──
def _send_slack(ticket_key, title, planner_json, developer_json, qa_json, report_url):
    developer = json.loads(developer_json)
    qa        = json.loads(qa_json)

    spec_quality = {a["id"]: a for a in qa.get("spec_quality", [])}
    label_map    = {"✅": "착수가능", "⚠️": "확인필요", "🔴": "착수불가"}
    stats        = {"착수가능": 0, "확인필요": 0, "착수불가": 0}
    for a in spec_quality.values():
        v = a.get("verdict")
        if v:
            stats[label_map.get(v, "확인필요")] += 1

    risk_items = [
        f"• [{a['id']}]: {a.get('missing_info', '')[:60]}"
        for a in spec_quality.values() if a.get("verdict") == "🔴"
    ]
    risk_block = "\n".join(risk_items[:3]) if risk_items else "없음"

    cqs = qa.get("critical_questions", [])
    p0_block = "\n".join(
        f"• {q['question'][:80]}"
        for q in cqs if q.get("priority") == "P0"
    )[:600] or "없음"

    high_risks = [r for r in developer.get("impl_risks", []) if r.get("complexity") == "High"]
    dev_block  = "\n".join(
        f"🔴 `High` {r['feature']} — {r['risk'][:50]}"
        for r in high_risks[:3]
    ) or "없음"

    text = f"""🤖 *기획서 품질 분석 완료*
티켓: *{ticket_key}* — {title}
분석일: {datetime.date.today()}

📊 *기획서 착수 가능성 판정*
✅ 착수가능: {stats['착수가능']}건　⚠️ 확인필요: {stats['확인필요']}건　🔴 착수불가: {stats['착수불가']}건

❓ *Critical Questions (P0 — 착수 전 기획자 확인 필요)*
{p0_block}

🔴 *착수불가 항목*
{risk_block}

🛠 *개발 구현 리스크 (High)*
{dev_block}

📄 *분석 보고서*: {report_url}"""

    app.client.chat_postMessage(channel=CHANNEL, text=text)


# ── Flask 라우트 ──
@flask_app.route("/slack/events", methods=["POST"])
def slack_events():
    return handler.handle(request)

@flask_app.route("/health")
def health():
    return "ok"

if __name__ == "__main__":
    port = int(os.environ.get("PORT", 3000))
    flask_app.run(host="0.0.0.0", port=port)
